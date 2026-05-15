import uuid
from io import BytesIO
from datetime import datetime, timezone
from pathlib import Path

from fastapi import HTTPException
from qdrant_client.http import models

from app.ai.embeddings.chunking import chunk_text
from app.core.config import settings
from app.services.vector_service import VectorService


class DocumentService:
    def __init__(self, vector_service=None, embedding_provider=None) -> None:
        self.vector_service = vector_service or VectorService(settings.QDRANT_DOCUMENTS_COLLECTION)
        self._embedding_provider = embedding_provider

    @property
    def embedding_provider(self):
        if self._embedding_provider is None:
            from app.ai.embeddings.local import FastEmbedProvider

            self._embedding_provider = FastEmbedProvider()
        return self._embedding_provider

    async def upload_and_ingest(
        self,
        user_id: uuid.UUID,
        file_name: str,
        content: str | None = None,
        file_bytes: bytes | None = None,
        file_type: str = "text/plain",
        scope: str = "global",
        chat_id: uuid.UUID | None = None,
        selected_model_id: str | None = None,
        selected_model_supports_vision: bool = False,
    ) -> dict:
        normalized_scope = "chat" if scope == "chat" else "global"
        if normalized_scope == "chat" and chat_id is None:
            raise HTTPException(status_code=400, detail="Chat-scoped documents require a chat id")
        raw_bytes = file_bytes if file_bytes is not None else (content or "").encode("utf-8")
        if not raw_bytes:
            raise HTTPException(status_code=400, detail="Uploaded document is empty")

        document_id = uuid.uuid4()
        safe_name = Path(file_name or f"{document_id}.txt").name
        object_path = f"{settings.MINIO_BUCKET}/{user_id}/{document_id}/{safe_name}"
        target = Path(settings.MINIO_DATA_PATH) / object_path
        target.parent.mkdir(parents=True, exist_ok=True)
        target.write_bytes(raw_bytes)

        is_image = file_type.startswith("image/") or Path(safe_name).suffix.lower() in {".png", ".jpg", ".jpeg", ".webp"}
        requires_multimodal = is_image and not selected_model_supports_vision
        indexable_content = extract_document_text(safe_name, file_type, raw_bytes, fallback=content)
        if is_image:
            indexable_content = (
                f"Image document: {safe_name}. "
                f"Uploaded for {'current chat' if normalized_scope == 'chat' else 'global knowledge library'}. "
                f"Selected model: {selected_model_id or 'not selected'}. "
            )
            if requires_multimodal:
                indexable_content += "Image understanding requires a multimodal or vision-capable model."

        if not indexable_content.strip():
            raise HTTPException(status_code=400, detail="No indexable text found in document")
        chunks = chunk_text(indexable_content)
        if not chunks:
            raise HTTPException(status_code=400, detail="No indexable text found in document")
        vectors = await self.embedding_provider.embed(chunks)
        uploaded_at = datetime.now(timezone.utc).isoformat()
        payloads = [
            {
                "source_type": "document",
                "user_id": str(user_id),
                "document_id": str(document_id),
                "source_id": str(document_id),
                "file_name": safe_name,
                "file_type": file_type,
                "file_size": len(raw_bytes),
                "scope": normalized_scope,
                "chat_id": str(chat_id) if chat_id else None,
                "status": "indexed",
                "requires_multimodal": requires_multimodal,
                "chunk_id": str(uuid.uuid4()),
                "chunk_index": index,
                "uploaded_date": uploaded_at,
                "minio_object_path": object_path,
                "title": safe_name,
                "text": chunk,
            }
            for index, chunk in enumerate(chunks)
        ]
        await self.vector_service.upsert(vectors, payloads)
        return {
            "document_id": document_id,
            "file_name": safe_name,
            "file_type": file_type,
            "uploaded_date": uploaded_at,
            "minio_object_path": object_path,
            "chunk_count": len(chunks),
            "file_size": len(raw_bytes),
            "scope": normalized_scope,
            "chat_id": chat_id,
            "status": "indexed",
            "requires_multimodal": requires_multimodal,
        }

    async def list_documents(self, user_id: uuid.UUID, scope: str | None = None, chat_id: uuid.UUID | None = None) -> list[dict]:
        must = [
            models.FieldCondition(key="user_id", match=models.MatchValue(value=str(user_id))),
            models.FieldCondition(key="source_type", match=models.MatchValue(value="document")),
        ]
        if scope:
            must.append(models.FieldCondition(key="scope", match=models.MatchValue(value=scope)))
        if chat_id:
            must.append(models.FieldCondition(key="chat_id", match=models.MatchValue(value=str(chat_id))))
        points = await self.vector_service.scroll(models.Filter(must=must), limit=500)
        documents: dict[str, dict] = {}
        for point in points:
            payload = point.payload or {}
            document_id = payload.get("document_id") or payload.get("source_id")
            if not document_id:
                continue
            entry = documents.setdefault(
                document_id,
                {
                    "document_id": document_id,
                    "file_name": payload.get("file_name") or payload.get("title") or "Untitled document",
                    "file_type": payload.get("file_type"),
                    "uploaded_date": payload.get("uploaded_date"),
                    "minio_object_path": payload.get("minio_object_path"),
                    "chunk_count": 0,
                    "scope": payload.get("scope") or "global",
                    "chat_id": payload.get("chat_id"),
                    "status": payload.get("status") or "indexed",
                    "requires_multimodal": bool(payload.get("requires_multimodal")),
                },
            )
            entry["chunk_count"] += 1
        return sorted(documents.values(), key=lambda item: item.get("uploaded_date") or "", reverse=True)

    async def delete_document(self, user_id: uuid.UUID, document_id: uuid.UUID) -> None:
        documents = await self.list_documents(user_id)
        document = next((item for item in documents if str(item["document_id"]) == str(document_id)), None)
        await self.vector_service.delete_source(user_id, "document", document_id)
        if document and document.get("minio_object_path"):
            target = Path(settings.MINIO_DATA_PATH) / document["minio_object_path"]
            if target.exists():
                target.unlink()

    async def update_scope(self, user_id: uuid.UUID, document_id: uuid.UUID, scope: str, chat_id: uuid.UUID | None) -> dict:
        normalized_scope = "chat" if scope == "chat" else "global"
        if normalized_scope == "chat" and chat_id is None:
            raise HTTPException(status_code=400, detail="Chat-specific documents require a chat id")
        documents = await self.list_documents(user_id)
        document = next((item for item in documents if str(item["document_id"]) == str(document_id)), None)
        if document is None:
            raise HTTPException(status_code=404, detail="Document not found")
        await self.vector_service.update_source_payload(
            user_id,
            "document",
            document_id,
            {"scope": normalized_scope, "chat_id": str(chat_id) if normalized_scope == "chat" else None},
        )
        document["scope"] = normalized_scope
        document["chat_id"] = str(chat_id) if normalized_scope == "chat" else None
        return document


def extract_document_text(file_name: str, file_type: str, data: bytes, fallback: str | None = None) -> str:
    extension = Path(file_name).suffix.lower()
    if fallback and fallback.strip():
        return fallback
    if file_type.startswith("text/") or extension in {".txt", ".md", ".markdown", ".csv", ".json"}:
        return data.decode("utf-8", errors="replace")
    if file_type == "application/pdf" or extension == ".pdf":
        return extract_pdf_text(data)
    if extension == ".docx" or file_type.endswith("wordprocessingml.document"):
        return extract_docx_text(data)
    if extension == ".pptx" or file_type.endswith("presentationml.presentation"):
        return extract_pptx_text(data)
    if file_type.startswith("image/") or extension in {".png", ".jpg", ".jpeg", ".webp"}:
        return f"Image document: {file_name}."
    return f"Uploaded binary document: {file_name}. Text extraction is not available for this file type yet."


def extract_pdf_text(data: bytes) -> str:
    try:
        from pypdf import PdfReader

        reader = PdfReader(BytesIO(data))
        return "\n\n".join(page.extract_text() or "" for page in reader.pages)
    except Exception:
        return ""


def extract_docx_text(data: bytes) -> str:
    try:
        from docx import Document

        document = Document(BytesIO(data))
        paragraphs = [paragraph.text for paragraph in document.paragraphs if paragraph.text.strip()]
        table_cells = [
            cell.text
            for table in document.tables
            for row in table.rows
            for cell in row.cells
            if cell.text.strip()
        ]
        return "\n".join([*paragraphs, *table_cells])
    except Exception:
        return ""


def extract_pptx_text(data: bytes) -> str:
    try:
        from pptx import Presentation

        presentation = Presentation(BytesIO(data))
        parts: list[str] = []
        for slide_index, slide in enumerate(presentation.slides, start=1):
            slide_text = []
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text.strip():
                    slide_text.append(shape.text)
            if slide_text:
                parts.append(f"Slide {slide_index}\n" + "\n".join(slide_text))
        return "\n\n".join(parts)
    except Exception:
        return ""
